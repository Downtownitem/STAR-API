import textract
import PyPDF2
from pptx import Presentation


class TextExtractor:
    def __init__(self, file_path):
        self.file_path = file_path

    def __extract_txt(self):
        with open(self.file_path, 'r') as file:
            return file.read()

    def __extract_doc(self):
        return textract.process(self.file_path).decode()

    def __extract_pdf(self):
        pdf_file_obj = open(self.file_path, 'rb')
        pdf_reader = PyPDF2.PdfReader(pdf_file_obj)
        num_pages = len(pdf_reader.pages)
        text = ""

        for page in range(num_pages):
            page_obj = pdf_reader.pages[page]
            text += page_obj.extract_text()

        pdf_file_obj.close()
        return text

    def __extract_ppt(self):
        presentation = Presentation(self.file_path)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text

        return text

    def extract(self):
        if self.file_path.endswith('.txt'):
            return self.__extract_txt()
        elif self.file_path.endswith('.doc') or self.file_path.endswith('.docx'):
            return self.__extract_doc()
        elif self.file_path.endswith('.pdf'):
            return self.__extract_pdf()
        elif self.file_path.endswith('.ppt') or self.file_path.endswith('.pptx'):
            return self.__extract_ppt()
        else:
            return "Tipo de archivo no soportado"
